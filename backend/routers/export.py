from fastapi import APIRouter, HTTPException, Header
from fastapi.responses import StreamingResponse
from docx import Document
from docx.shared import Pt, RGBColor
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.units import cm
from pptx import Presentation
from pptx.util import Inches, Pt as PPTXPt
from pptx.dml.color import RGBColor as PPTXColor
from supabase import create_client
import os, json, io

router = APIRouter()

def get_sb():
    return create_client(os.environ["SUPABASE_URL"], os.environ["SUPABASE_SERVICE_KEY"])

async def get_user_and_doc(doc_id: str, authorization: str):
    if not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401)
    token = authorization.split(" ")[1]
    sb = get_sb()
    user = sb.auth.get_user(token)
    user_id = user.user.id
    result = sb.table("documents").select("*").eq("id", doc_id).eq("user_id", user_id).single().execute()
    if not result.data:
        raise HTTPException(status_code=404)
    return user_id, result.data

@router.get("/{doc_id}/docx")
async def export_docx(doc_id: str, authorization: str = Header(...)):
    user_id, doc = await get_user_and_doc(doc_id, authorization)
    concepts = json.loads(doc.get("concepts_json") or "[]")
    kg = json.loads(doc.get("kg_json") or "{}")

    document = Document()
    
    # Başlık
    title = document.add_heading(doc["filename"].replace(".pdf", ""), 0)
    title.runs[0].font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
    
    document.add_paragraph(f"Sayfa sayısı: {doc['page_count']} | Kavram sayısı: {len(concepts)}")
    document.add_paragraph()

    # Kavramlar
    document.add_heading("Tespit Edilen Kavramlar", level=1)
    for c in concepts:
        p = document.add_paragraph()
        run = p.add_run(f"• {c.get('kavram', '')}: ")
        run.bold = True
        p.add_run(c.get("aciklama", ""))

    # KG özeti
    if kg.get("genel_not"):
        document.add_heading("Genel Değerlendirme", level=1)
        document.add_paragraph(kg["genel_not"])

    buf = io.BytesIO()
    document.save(buf)
    buf.seek(0)
    
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f"attachment; filename=bilgiharitasi_{doc_id[:8]}.docx"}
    )

@router.get("/{doc_id}/pdf-report")
async def export_pdf_report(doc_id: str, authorization: str = Header(...)):
    user_id, doc = await get_user_and_doc(doc_id, authorization)
    concepts = json.loads(doc.get("concepts_json") or "[]")

    buf = io.BytesIO()
    pdf = SimpleDocTemplate(buf, pagesize=A4, leftMargin=2*cm, rightMargin=2*cm)
    styles = getSampleStyleSheet()
    
    story = []
    title_style = ParagraphStyle("title", parent=styles["Title"], fontSize=20, textColor="#1a1a2e")
    story.append(Paragraph(doc["filename"].replace(".pdf", ""), title_style))
    story.append(Spacer(1, 0.5*cm))
    story.append(Paragraph(f"<b>BilgiHaritası Analiz Raporu</b>", styles["Normal"]))
    story.append(Spacer(1, 1*cm))
    
    story.append(Paragraph("<b>Tespit Edilen Kavramlar</b>", styles["Heading2"]))
    for c in concepts:
        story.append(Paragraph(f"• <b>{c.get('kavram','')}</b>: {c.get('aciklama','')}", styles["Normal"]))
        story.append(Spacer(1, 0.2*cm))
    
    pdf.build(story)
    buf.seek(0)
    
    return StreamingResponse(
        buf,
        media_type="application/pdf",
        headers={"Content-Disposition": f"attachment; filename=bilgiharitasi_{doc_id[:8]}.pdf"}
    )

@router.get("/{doc_id}/pptx")
async def export_pptx(doc_id: str, authorization: str = Header(...)):
    user_id, doc = await get_user_and_doc(doc_id, authorization)
    concepts = json.loads(doc.get("concepts_json") or "[]")

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]

    # Kapak slaytı
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = PPTXColor(0x1a, 0x1a, 0x2e)
    
    tf = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf.text_frame.text = doc["filename"].replace(".pdf", "")
    tf.text_frame.paragraphs[0].runs[0].font.size = PPTXPt(40)
    tf.text_frame.paragraphs[0].runs[0].font.color.rgb = PPTXColor(0xff, 0xff, 0xff)
    
    sub = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(1))
    sub.text_frame.text = "BilgiHaritası Analiz Raporu"
    sub.text_frame.paragraphs[0].runs[0].font.size = PPTXPt(20)
    sub.text_frame.paragraphs[0].runs[0].font.color.rgb = PPTXColor(0x91, 0x9b, 0xff)

    # Her 4 kavram için bir slayt
    for i in range(0, min(len(concepts), 20), 4):
        batch = concepts[i:i+4]
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = PPTXColor(0x0d, 0x11, 0x17)
        
        for j, c in enumerate(batch):
            x = Inches(0.5 + (j % 2) * 6.4)
            y = Inches(0.5 + (j // 2) * 3.3)
            box = slide.shapes.add_textbox(x, y, Inches(6), Inches(3))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = c.get("kavram", "")
            p.runs[0].font.size = PPTXPt(18)
            p.runs[0].font.bold = True
            p.runs[0].font.color.rgb = PPTXColor(0x5b, 0x66, 0xf1)
            
            tf.add_paragraph()
            p2 = tf.paragraphs[1]
            p2.text = c.get("aciklama", "")
            p2.runs[0].font.size = PPTXPt(12)
            p2.runs[0].font.color.rgb = PPTXColor(0xcc, 0xcc, 0xdd)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename=bilgiharitasi_{doc_id[:8]}.pptx"}
    )
